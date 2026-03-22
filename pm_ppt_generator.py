#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PM PPT Generator
마크다운 산출물을 Gemini AI로 구조화한 뒤 python-pptx로 .pptx를 생성합니다.
스킬별 전용 슬라이드 레이아웃을 지원합니다.

사용법: python pm_ppt_generator.py <markdown_file_path>
"""

import sys
import os
import re
import json
import io
from pathlib import Path
from datetime import datetime

# .env 자동 로드
_ENV_FILE = Path(r"C:\Users\psh93\OneDrive\Desktop\Claude\.scripts\.env")
if _ENV_FILE.exists():
    for _line in _ENV_FILE.read_text(encoding="utf-8").splitlines():
        _line = _line.strip()
        if _line and not _line.startswith("#") and "=" in _line:
            _k, _v = _line.split("=", 1)
            os.environ.setdefault(_k.strip(), _v.strip())

OUTPUT_DIR = Path(r"C:\Users\psh93\OneDrive\Desktop\Claude\05_PM_Outputs")

# ── 스킬별 Gemini 프롬프트 ────────────────────────────────────────────────────

SKILL_PROMPTS = {

"brainstorm-okrs": """
당신은 OKR 전문 PM입니다. 아래 OKR 산출물 마크다운을 분석해
발표용 PPT 슬라이드 JSON 배열을 생성하세요.

슬라이드 구성 규칙 (반드시 이 순서와 타입을 따르세요):

1. type=goal (1장)
   - 제품의 핵심 방향성과 정체성을 담은 Goal 슬라이드
   - goal: 한 문장으로 표현한 제품의 궁극적 목표
   - direction: 제품이 나아가는 방향 설명 (2-3문장, 구체적으로)
   - tagline: 팀을 하나로 묶는 짧은 슬로건 (15자 이내)

2. type=objective (Objective마다 1장씩, 최대 3장)
   - 각 Objective별 슬라이드
   - index: 번호 (1, 2, 3)
   - objective: Objective 전체 문장
   - why: 이 Objective가 Goal 달성에 필요한 이유 (구체적, 2문장)
   - linked_krs: 이 Objective에 연결된 KR 미리보기 (2-3개, 짧게)

3. type=key_results (1장, 전체 KR 통합)
   - 모든 Objective의 KR을 하나의 테이블로 정리
   - items: 각 KR 항목
     - objective: 어떤 Objective의 KR인지 (짧게)
     - kr: KR 설명 (짧게)
     - baseline: 현재 기준값 (없으면 "측정 필요")
     - target: 목표 수치 (구체적인 숫자 또는 %)

4. type=kpi (1장)
   - KR을 뒷받침하는 정량 지표와 산출 수식
   - items: KPI 항목
     - kpi: 지표명
     - formula: 산출 수식 (예: (금월 MAU - 전월 MAU) / 전월 MAU × 100)
     - target: 목표치
     - source: 데이터 출처 또는 측정 방법

JSON만 반환하고 코드블록(```)은 제거하세요.

출력 형식:
[
  {
    "type": "goal",
    "title": "Nova Platform 2026 Q2 Goal",
    "goal": "...",
    "direction": "...",
    "tagline": "..."
  },
  {
    "type": "objective",
    "index": 1,
    "title": "Objective 1: ...",
    "objective": "...",
    "why": "...",
    "linked_krs": ["KR1 미리보기", "KR2 미리보기"]
  },
  {
    "type": "key_results",
    "title": "Key Results 전체 현황",
    "items": [
      {"objective": "...", "kr": "...", "baseline": "...", "target": "..."}
    ]
  },
  {
    "type": "kpi",
    "title": "KPI & 달성도 산출 수식",
    "items": [
      {"kpi": "...", "formula": "...", "target": "...", "source": "..."}
    ]
  }
]

분석할 OKR 산출물:
{content}
""",

}  # 다른 스킬은 GENERIC_PROMPT 사용

GENERIC_PROMPT = """
아래 PM 산출물을 분석해 발표용 슬라이드 JSON 배열을 생성하세요.

규칙:
- 각 슬라이드: type="generic", title, subtitle(선택), bullets(최대 5개)
- bullets는 짧고 명확한 한국어 (한 줄 30자 이내)
- 슬라이드 수: 5~8장 (커버·목차 제외)
- JSON만 반환, 코드블록(```) 제거

형식:
[
  {{"type": "generic", "title": "...", "subtitle": "...", "bullets": ["..."]}}
]

산출물:
{content}
"""


# ── 프론트매터 파싱 ───────────────────────────────────────────────────────────

def parse_frontmatter(md_text: str) -> dict:
    meta = {"skill": "", "context": "", "date": ""}
    m = re.match(r"^---\s*\n(.+?)\n---", md_text, re.DOTALL)
    if not m:
        return meta
    for line in m.group(1).splitlines():
        if ":" in line:
            k, v = line.split(":", 1)
            key = k.strip().lower()
            if key in meta:
                meta[key] = v.strip()
    return meta


# ── Gemini API 호출 ───────────────────────────────────────────────────────────

def call_gemini(md_text: str, skill: str) -> list[dict]:
    try:
        from google import genai
    except ImportError:
        raise RuntimeError("google-genai 패키지 없음. pip install google-genai")

    api_key = os.environ.get("GEMINI_API_KEY", "")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY 미설정. .scripts/.env 확인")

    template = SKILL_PROMPTS.get(skill, GENERIC_PROMPT)
    prompt = template.replace("{content}", md_text)

    client = genai.Client(api_key=api_key)
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
    )
    raw = response.text.strip()
    raw = re.sub(r"```json\s*", "", raw)
    raw = re.sub(r"```\s*", "", raw)
    return json.loads(raw)


# ── python-pptx 공통 유틸 ────────────────────────────────────────────────────

def _setup_prs():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs, prs.slide_layouts[6]  # blank layout


def _palette():
    from pptx.dml.color import RGBColor
    return {
        "bg":       RGBColor(0x1A, 0x1A, 0x2E),
        "accent":   RGBColor(0xE9, 0x4F, 0x37),
        "text":     RGBColor(0xF0, 0xF0, 0xF0),
        "bullet":   RGBColor(0xA8, 0xD8, 0xEA),
        "dim":      RGBColor(0x88, 0x88, 0x88),
        "green":    RGBColor(0x6B, 0xC5, 0x8A),
        "yellow":   RGBColor(0xF7, 0xC5, 0x48),
        "sub":      RGBColor(0xCC, 0xCC, 0xCC),
        "table_hd": RGBColor(0x2A, 0x2A, 0x4A),
        "table_row":RGBColor(0x22, 0x22, 0x3A),
    }


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _bar(slide, top, C):
    from pptx.util import Inches
    b = slide.shapes.add_shape(1, Inches(0), Inches(top), Inches(13.33), Inches(0.07))
    b.fill.solid()
    b.fill.fore_color.rgb = C
    b.line.fill.background()


def _tb(slide, l, t, w, h, text, pt, bold=False, italic=False, color=None,
        wrap=True, align=None):
    """텍스트박스 추가 헬퍼."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(pt)
    p.font.bold = bold
    p.font.italic = italic
    if color:
        p.font.color.rgb = color
    if align:
        p.alignment = align
    return tf


def _slide_num(slide, num_str, C):
    from pptx.util import Inches, Pt
    tx = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.text = num_str
    p.font.size = Pt(11)
    p.font.color.rgb = C


# ── 슬라이드 렌더러 ──────────────────────────────────────────────────────────

def _render_cover(prs, blank, meta, C):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(blank)
    _bg(slide, C["bg"])
    _bar(slide, 0, C["accent"])
    _bar(slide, 7.43, C["accent"])

    ctx      = meta.get("context", "")
    skill    = meta.get("skill", "").upper()
    date_str = meta.get("date", datetime.today().strftime("%Y-%m-%d"))

    _tb(slide, 1.2, 1.8, 10.9, 2.2, ctx or skill, 38, bold=True, color=C["text"])
    _tb(slide, 1.2, 4.3, 10.9, 0.7,
        f"{skill}  ·  {date_str}", 17, color=C["sub"])


def _render_toc(prs, blank, slides, C):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(blank)
    _bg(slide, C["bg"])
    _bar(slide, 1.25, C["accent"])

    _tb(slide, 0.6, 0.3, 12, 0.9, "목차", 30, bold=True, color=C["accent"])

    tx = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.5), Inches(5.6))
    tf = tx.text_frame
    tf.word_wrap = True

    TYPE_ICON = {"goal": "◈", "objective": "▷", "key_results": "◆", "kpi": "⊕", "generic": "•"}
    for idx, s in enumerate(slides):
        t = s.get("type", "generic")
        icon = TYPE_ICON.get(t, "•")
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"  {icon}  {s.get('title', '')}"
        p.font.size = Pt(19)
        p.font.color.rgb = C["bullet"]
        p.space_before = Pt(5)

    _slide_num(slide, "—", C["dim"])


def _render_goal(prs, blank, data, slide_num, C):
    """Goal 슬라이드: 방향성 + tagline 전면 배치."""
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(blank)
    _bg(slide, C["bg"])
    _bar(slide, 0, C["accent"])

    _tb(slide, 0.6, 0.18, 12, 0.7, data.get("title", "Goal"), 13,
        bold=True, color=C["accent"])

    # Goal 핵심 문장 — 크고 중앙
    goal_text = data.get("goal", "")
    _tb(slide, 0.8, 0.9, 11.7, 1.6, goal_text, 28, bold=True, color=C["text"])

    # 구분선
    _bar(slide, 2.6, C["dim"])

    # 방향성 설명
    direction = data.get("direction", "")
    _tb(slide, 0.8, 2.75, 11.5, 2.4, direction, 19, color=C["bullet"], italic=False)

    # Tagline 박스
    tagline = data.get("tagline", "")
    if tagline:
        tag_box = slide.shapes.add_shape(1, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.85))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = C["accent"]
        tag_box.line.fill.background()
        from pptx.util import Pt as P
        from pptx.enum.text import PP_ALIGN
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = f'" {tagline} "'
        p_tag.font.size = P(22)
        p_tag.font.bold = True
        p_tag.font.color.rgb = C["text"]
        p_tag.alignment = PP_ALIGN.CENTER

    _bar(slide, 7.43, C["accent"])
    _slide_num(slide, str(slide_num), C["dim"])


def _render_objective(prs, blank, data, slide_num, C):
    """Objective 슬라이드: O 전문 + Why + KR 미리보기."""
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(blank)
    _bg(slide, C["bg"])

    idx   = data.get("index", "")
    title = data.get("title", f"Objective {idx}")
    _tb(slide, 0.6, 0.25, 12, 0.75, title, 26, bold=True, color=C["accent"])
    _bar(slide, 1.1, C["accent"])

    # O 전문
    obj_text = data.get("objective", "")
    _tb(slide, 0.7, 1.25, 11.9, 1.3, obj_text, 22, bold=True, color=C["text"])

    # Why 레이블 + 내용
    _tb(slide, 0.7, 2.65, 1.5, 0.45, "WHY", 13, bold=True, color=C["accent"])
    why_text = data.get("why", "")
    _tb(slide, 0.7, 3.05, 11.8, 1.5, why_text, 18, color=C["bullet"])

    # KR 미리보기
    _bar(slide, 4.7, C["table_hd"])
    _tb(slide, 0.7, 4.8, 3.0, 0.4, "연결된 Key Results", 13, bold=True, color=C["yellow"])
    krs = data.get("linked_krs", [])
    tx = slide.shapes.add_textbox(Inches(0.9), Inches(5.25), Inches(11.5), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, kr in enumerate(krs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  ▸  {kr}"
        p.font.size = Pt(17)
        p.font.color.rgb = C["green"]
        p.space_before = Pt(4)

    _slide_num(slide, str(slide_num), C["dim"])


def _render_key_results(prs, blank, data, slide_num, C):
    """Key Results 슬라이드: 4열 테이블 (Objective / KR / Baseline / Target)."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(blank)
    _bg(slide, C["bg"])

    title = data.get("title", "Key Results")
    _tb(slide, 0.6, 0.25, 12, 0.75, title, 26, bold=True, color=C["accent"])
    _bar(slide, 1.1, C["accent"])

    items = data.get("items", [])
    rows  = len(items) + 1  # 헤더 포함
    cols  = 4
    col_w = [3.2, 4.5, 2.3, 2.3]  # 합 = 12.3

    tbl = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.3),
        Inches(sum(col_w)), Inches(min(5.8, 0.55 + rows * 0.72))
    ).table

    # 열 너비
    for ci, w in enumerate(col_w):
        tbl.columns[ci].width = Inches(w)

    # 헤더
    headers = ["Objective", "Key Result", "Baseline", "Target"]
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["accent"]
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C["text"]
        p.alignment = PP_ALIGN.CENTER

    # 데이터 행
    for ri, item in enumerate(items, start=1):
        row_vals = [
            item.get("objective", ""),
            item.get("kr", ""),
            item.get("baseline", "측정 필요"),
            item.get("target", ""),
        ]
        row_color = C["table_row"] if ri % 2 == 0 else C["table_hd"]
        for ci, val in enumerate(row_vals):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(13)
            txt_color = C["green"] if ci == 3 else C["text"]
            p.font.color.rgb = txt_color
            if ci == 3:
                p.alignment = PP_ALIGN.CENTER

    _slide_num(slide, str(slide_num), C["dim"])


def _render_kpi(prs, blank, data, slide_num, C):
    """KPI 슬라이드: 4열 테이블 (KPI / 산출 수식 / Target / 출처)."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(blank)
    _bg(slide, C["bg"])

    title = data.get("title", "KPI & 달성도 산출 수식")
    _tb(slide, 0.6, 0.25, 12, 0.75, title, 26, bold=True, color=C["accent"])
    _bar(slide, 1.1, C["accent"])

    items = data.get("items", [])
    rows  = len(items) + 1
    cols  = 4
    col_w = [2.4, 5.2, 2.0, 2.7]

    tbl = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.3),
        Inches(sum(col_w)), Inches(min(5.8, 0.55 + rows * 0.80))
    ).table

    for ci, w in enumerate(col_w):
        tbl.columns[ci].width = Inches(w)

    headers = ["KPI 지표", "산출 수식", "목표치", "측정 출처"]
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["accent"]
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C["text"]
        p.alignment = PP_ALIGN.CENTER

    for ri, item in enumerate(items, start=1):
        row_vals = [
            item.get("kpi", ""),
            item.get("formula", ""),
            item.get("target", ""),
            item.get("source", ""),
        ]
        row_color = C["table_row"] if ri % 2 == 0 else C["table_hd"]
        for ci, val in enumerate(row_vals):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(12)
            txt_color = C["yellow"] if ci == 2 else (C["bullet"] if ci == 1 else C["text"])
            p.font.color.rgb = txt_color

    _slide_num(slide, str(slide_num), C["dim"])


def _render_generic(prs, blank, data, slide_num, C):
    """일반 슬라이드 렌더러 (비OKR 스킬용)."""
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(blank)
    _bg(slide, C["bg"])

    _tb(slide, 0.6, 0.3, 12, 1.0, data.get("title", ""), 30,
        bold=True, color=C["accent"])

    sub = data.get("subtitle", "")
    top = 1.5
    if sub:
        _tb(slide, 0.6, top, 12, 0.55, sub, 17, italic=True, color=C["sub"])
        top += 0.65

    bullets = data.get("bullets", [])
    if bullets:
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.5), Inches(5.5))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(20)
            p.font.color.rgb = C["bullet"]
            p.space_before = Pt(8)

    _slide_num(slide, str(slide_num), C["dim"])


# ── 메인 빌드 ────────────────────────────────────────────────────────────────

def build_pptx(slides: list[dict], output_path: Path, meta: dict) -> None:
    prs, blank = _setup_prs()
    C = _palette()

    _render_cover(prs, blank, meta, C)
    _render_toc(prs, blank, slides, C)

    RENDERERS = {
        "goal":        _render_goal,
        "objective":   _render_objective,
        "key_results": _render_key_results,
        "kpi":         _render_kpi,
    }

    for i, slide_data in enumerate(slides):
        stype    = slide_data.get("type", "generic")
        renderer = RENDERERS.get(stype, _render_generic)
        renderer(prs, blank, slide_data, i + 1, C)

    prs.save(str(output_path))


# ── 인덱스 README 갱신 ───────────────────────────────────────────────────────

def update_index(output_dir: Path) -> None:
    entries = []
    for md in sorted(output_dir.glob("*.md"), reverse=True):
        if md.name == "README.md":
            continue
        meta    = parse_frontmatter(md.read_text(encoding="utf-8", errors="ignore"))
        skill   = meta.get("skill") or md.stem
        context = meta.get("context") or ""
        date    = meta.get("date") or ""
        pptx    = md.with_suffix(".pptx")
        pptx_lk = f" · [PPT]({pptx.name})" if pptx.exists() else ""
        entries.append(f"| {date} | {skill} | {context} | [MD]({md.name}){pptx_lk} |")

    lines = [
        "# PM Outputs Index\n",
        f"_마지막 갱신: {datetime.today().strftime('%Y-%m-%d %H:%M')}_\n\n",
        "| 날짜 | 스킬 | 컨텍스트 | 파일 |",
        "|---|---|---|---|",
    ] + entries
    readme = output_dir / "README.md"
    readme.write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"      인덱스 갱신: {readme}")


# ── 메인 ─────────────────────────────────────────────────────────────────────

def main():
    sys.stdout = __import__("io").TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    if len(sys.argv) < 2:
        print("사용법: python pm_ppt_generator.py <markdown_file_path>")
        sys.exit(1)

    md_path = Path(sys.argv[1])
    if not md_path.exists():
        print(f"[오류] 파일 없음: {md_path}")
        sys.exit(1)

    print(f"\n[1/4] 마크다운 읽는 중: {md_path.name}")
    md_text = md_path.read_text(encoding="utf-8")
    meta    = parse_frontmatter(md_text)
    skill   = meta.get("skill", "")
    print(f"      skill={skill}  date={meta['date']}")

    print("[2/4] Gemini AI로 슬라이드 구조화 중...")
    slides = call_gemini(md_text, skill)
    types  = [s.get("type", "?") for s in slides]
    print(f"      {len(slides)}장 생성: {types}")

    stem        = md_path.stem
    output_path = OUTPUT_DIR / f"{stem}.pptx"
    print(f"[3/4] .pptx 생성 중: {output_path.name}")
    build_pptx(slides, output_path, meta)

    print("[4/4] 인덱스 갱신 중...")
    update_index(OUTPUT_DIR)

    total = len(slides) + 2
    print()
    print("━" * 56)
    print(f"✓ PPT 저장: {output_path}")
    print(f"  구성: 커버 + 목차 + 본문 {len(slides)}장 = 총 {total}장")
    print(f"  슬라이드 타입: {types}")
    print("━" * 56)


if __name__ == "__main__":
    main()
